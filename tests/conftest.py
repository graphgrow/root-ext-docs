"""Fixture builders — every test document is generated in-test (no
binary fixtures in the repo): PDFs are hand-assembled with a computed
xref (pdfium is the arbiter that they're valid), images via Pillow,
office files via their own libraries, archives via the stdlib."""

from __future__ import annotations

from pathlib import Path

import pytest


def build_pdf(path: Path, page_texts: list[str | None]) -> Path:
    """Assemble a minimal valid PDF: one page per entry; None makes a
    page with NO text layer (a stand-in for a scan/drawing)."""
    objects: list[bytes] = []
    page_count = len(page_texts)
    font_obj = 3 + 2 * page_count

    kids = " ".join(f"{3 + 2 * i} 0 R" for i in range(page_count))
    objects.append(b"<< /Type /Catalog /Pages 2 0 R >>")
    objects.append(
        f"<< /Type /Pages /Kids [{kids}] /Count {page_count} >>".encode()
    )
    for index, text in enumerate(page_texts):
        content_obj = 4 + 2 * index
        objects.append(
            (
                f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
                f"/Contents {content_obj} 0 R "
                f"/Resources << /Font << /F1 {font_obj} 0 R >> >> >>"
            ).encode()
        )
        if text is None:
            stream = b""
        else:
            escaped = text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
            stream = f"BT /F1 24 Tf 72 720 Td ({escaped}) Tj ET".encode()
        objects.append(
            b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n"
            + stream + b"\nendstream"
        )
    objects.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    out = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{number} 0 obj\n".encode() + body + b"\nendobj\n"
    xref_at = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets[1:]:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_at}\n%%EOF\n"
    ).encode()
    path.write_bytes(bytes(out))
    return path


@pytest.fixture
def pdf_two_pages(tmp_path: Path) -> Path:
    """Page 1 has a text layer; page 2 is empty (the 'scan')."""
    return build_pdf(tmp_path / "sheet.pdf", ["Torque spec 42 Nm", None])


@pytest.fixture
def sample_image(tmp_path: Path) -> Path:
    from PIL import Image

    path = tmp_path / "plot.png"
    Image.new("RGB", (320, 200), (240, 240, 236)).save(path)
    return path


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    from docx import Document

    path = tmp_path / "notes.docx"
    document = Document()
    document.add_heading("Interface Control", level=1)
    document.add_paragraph("The bracket mounts with four M4 bolts.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Bolt"
    table.cell(0, 1).text = "Torque"
    table.cell(1, 0).text = "M4"
    table.cell(1, 1).text = "2.9 Nm"
    document.save(str(path))
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation

    path = tmp_path / "review.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Design Review"
    slide.placeholders[1].text = "Mass budget holds at 1.2 kg"
    slide.notes_slide.notes_text_frame.text = "Confirm with thermal team"
    presentation.save(str(path))
    return path


@pytest.fixture
def sample_zip(tmp_path: Path) -> Path:
    import zipfile

    path = tmp_path / "vendor_drop.zip"
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("readme.txt", "hello")
        archive.writestr("cad/bracket.step", "ISO-10303-21;")
    return path
