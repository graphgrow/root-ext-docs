"""Office extraction: DOCX and PPTX to structured markdown.

Richer than a raw text run: headings keep their level, tables become
markdown tables, slides keep title/body/notes separation — so the model
can cite "the table under §3" instead of fishing in a text soup.
"""

from __future__ import annotations

from . import files


def _table_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    def line(cells: list[str]) -> str:
        return "| " + " | ".join(c.replace("|", "\\|").replace("\n", " ") for c in cells) + " |"
    out = [line(rows[0]), line(["---"] * len(rows[0]))]
    out.extend(line(row) for row in rows[1:])
    return "\n".join(out)


def docx_extract(path: str) -> dict:
    """Document-order extraction: headings, paragraphs, tables."""
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    resolved = files.checked_path(path, files.DOCX_SUFFIXES)
    try:
        document = Document(str(resolved))
    except Exception as exc:
        raise ValueError(f"not a readable DOCX: {exc}") from exc

    blocks: list[str] = []
    tables = 0
    # iter_inner_content walks paragraphs and tables in document order
    # (python-docx >= 1.1); fall back to the flat lists if absent.
    items = (
        document.iter_inner_content()
        if hasattr(document, "iter_inner_content")
        else [*document.paragraphs, *document.tables]
    )
    for item in items:
        if isinstance(item, Paragraph):
            text = item.text.strip()
            if not text:
                continue
            style = (item.style.name or "") if item.style else ""
            if style.startswith("Heading"):
                try:
                    level = int(style.split()[-1])
                except ValueError:
                    level = 2
                blocks.append(f"{'#' * min(level, 6)} {text}")
            elif style == "Title":
                blocks.append(f"# {text}")
            else:
                blocks.append(text)
        elif isinstance(item, Table):
            tables += 1
            blocks.append(
                _table_markdown([[cell.text for cell in row.cells] for row in item.rows])
            )
    text, truncated = files.budget_text("\n\n".join(blocks))
    return {
        "source": str(resolved),
        "tables": tables,
        "text": text,
        "truncated": truncated,
    }


def pptx_extract(path: str) -> dict:
    """Slide-by-slide extraction: title, body text, tables, notes."""
    from pptx import Presentation

    resolved = files.checked_path(path, files.PPTX_SUFFIXES)
    try:
        presentation = Presentation(str(resolved))
    except Exception as exc:
        raise ValueError(f"not a readable PPTX: {exc}") from exc

    blocks: list[str] = []
    for number, slide in enumerate(presentation.slides, start=1):
        title = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip()
        blocks.append(f"## Slide {number}" + (f" — {title}" if title else ""))
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                body = shape.text_frame.text.strip()
                if body:
                    blocks.append(body)
            if getattr(shape, "has_table", False) and shape.has_table:
                blocks.append(
                    _table_markdown(
                        [[cell.text for cell in row.cells] for row in shape.table.rows]
                    )
                )
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                blocks.append(f"> Notes: {notes}")
    text, truncated = files.budget_text("\n\n".join(blocks))
    return {
        "source": str(resolved),
        "slides": len(presentation.slides),
        "text": text,
        "truncated": truncated,
    }
